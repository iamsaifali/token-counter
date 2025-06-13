import streamlit as st
import tiktoken
import pandas as pd
import numpy as np
from typing import Dict, List, Tuple, Any
import json
import io
import base64
from datetime import datetime
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
import os

# File processing imports
import docx
import PyPDF2
import pdfplumber
import openpyxl
import xlrd
from pptx import Presentation
import chardet
from PIL import Image

# Model management imports
try:
    from model_management_ui import show_model_management_ui
except ImportError:
    def show_model_management_ui(model_manager):
        st.error("Model Management UI not available. Please check installation.")

# Configure Streamlit page
st.set_page_config(
    page_title="OpenAI Token Counter & Cost Calculator",
    page_icon="🧮",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS for better styling
st.markdown("""
<style>
    .main-header {
        font-size: 3rem;
        font-weight: bold;
        text-align: center;
        color: #1E90FF;
        margin-bottom: 2rem;
    }
    .subheader {
        font-size: 1.5rem;
        font-weight: bold;
        color: #4682B4;
        margin-top: 2rem;
        margin-bottom: 1rem;
    }
    .metric-card {
        background-color: #F0F8FF;
        padding: 1rem;
        border-radius: 10px;
        border-left: 5px solid #1E90FF;
        margin: 1rem 0;
    }
    .cost-breakdown {
        background-color: #F5F5F5;
        padding: 1.5rem;
        border-radius: 10px;
        margin: 1rem 0;
    }
    .model-info {
        background-color: #E6F3FF;
        padding: 1rem;
        border-radius: 8px;
        margin: 0.5rem 0;
    }
    .warning-box {
        background-color: #FFF3CD;
        border: 1px solid #FFEAA7;
        padding: 1rem;
        border-radius: 8px;
        margin: 1rem 0;
    }
    .category-header {
        background-color: #4682B4;
        color: white;
        padding: 0.5rem 1rem;
        border-radius: 5px;
        margin: 1rem 0;
        font-weight: bold;
    }
</style>
""", unsafe_allow_html=True)

class OpenAIModels:
    """OpenAI model pricing and information with multi-category support."""
    
    def __init__(self):
        # Import and use the new model manager
        try:
            from model_manager import ModelManager
            self.model_manager = ModelManager()
            self.models = self.model_manager.get_all_models()
            self.image_generation_models = self.model_manager.get_models_by_category('image_generation_models')
            self.web_search_models = self.model_manager.get_models_by_category('web_search_models')
            self.built_in_tools = self.model_manager.get_models_by_category('built_in_tools')
            
            # Apply pricing types based on category
            self.apply_pricing_types()
            
        except Exception as e:
            st.error(f"Error initializing model manager: {e}")
            # Only use emergency fallback if everything fails
            self.load_emergency_fallback()
    
    def apply_pricing_types(self):
        """Apply pricing types based on model categories"""
        for model_id, model_info in self.models.items():
            category = model_info.get('category', '').lower()
            
            # Set pricing_type if not already set
            if 'pricing_type' not in model_info:
                if any(audio_term in category for audio_term in ['audio', 'realtime']):
                    if 'realtime' in category:
                        model_info['pricing_type'] = 'realtime_audio'
                    else:
                        model_info['pricing_type'] = 'audio_tokens'
                elif any(speech_term in category for speech_term in ['transcription', 'whisper', 'tts']):
                    model_info['pricing_type'] = 'per_minute'
                elif 'embedding' in category:
                    model_info['pricing_type'] = 'embeddings'
                elif 'fine' in category:
                    model_info['pricing_type'] = 'fine_tuning'
                elif 'moderation' in category:
                    model_info['pricing_type'] = 'free'
                else:
                    model_info['pricing_type'] = 'standard'
    
    def load_emergency_fallback(self):
        """Emergency fallback with minimal essential models"""
        st.warning("⚠️ Using emergency fallback models. Please check your pricing data files.")
        self.models = {
            "gpt-4o-2024-11-20": {
                "name": "GPT-4o",
                "description": "High-intelligence multimodal model",
                "input_cost": 2.50,
                "cached_input_cost": 1.25,
                "output_cost": 10.00,
                "context_window": 128000,
                "category": "Flagship Models",
                "has_cached": True,
                "pricing_type": "standard"
            },
            "gpt-4o-mini-2024-07-18": {
                "name": "GPT-4o Mini",
                "description": "Fast and affordable multimodal model",
                "input_cost": 0.15,
                "cached_input_cost": 0.075,
                "output_cost": 0.60,
                "context_window": 128000,
                "category": "Small Models",
                "has_cached": True,
                "pricing_type": "standard"
            }
        }
        self.image_generation_models = {}
        self.web_search_models = {}
        self.built_in_tools = {}
    
    @staticmethod
    def get_all_models() -> Dict[str, Dict[str, Any]]:
        """Get all models from the instance."""
        instance = OpenAIModels()
        return instance.models
    
    def calculate_cost_by_category(self, model_id: str, **kwargs) -> Tuple[float, Dict[str, Any]]:
        """Calculate cost based on model category and pricing type."""
        
        # Check if it's a regular text model
        if model_id in self.models:
            model = self.models[model_id]
            pricing_type = model.get('pricing_type', 'standard')
            
            if pricing_type == 'audio_tokens':
                return self.calculate_audio_model_cost(model_id, **kwargs)
            elif pricing_type == 'realtime_audio':
                return self.calculate_realtime_audio_cost(model_id, **kwargs)
            elif pricing_type == 'per_minute':
                return self.calculate_per_minute_cost(model_id, **kwargs)
            elif pricing_type == 'per_character':
                return self.calculate_per_character_cost(model_id, **kwargs)
            elif pricing_type == 'fine_tuning':
                return self.calculate_fine_tuning_cost(model_id, **kwargs)
            elif pricing_type == 'embeddings':
                return self.calculate_embeddings_cost(model_id, **kwargs)
            elif model.get('is_free'):
                return self.calculate_free_model_cost(model_id, **kwargs)
            else:
                # Standard text models (including o1, GPT-4o, GPT-3.5, etc.)
                return self.calculate_text_model_cost(model_id, **kwargs)
        
        # Check if it's an image generation model
        elif model_id in self.image_generation_models:
            return self.calculate_image_generation_cost(model_id, **kwargs)
        
        # Check if it's a web search model
        elif model_id in self.web_search_models:
            return self.calculate_web_search_cost(model_id, **kwargs)
        
        # Check if it's a built-in tool
        elif model_id in self.built_in_tools:
            return self.calculate_tool_cost(model_id, **kwargs)
        
        else:
            raise ValueError(f"Model {model_id} not found in any category")
    
    def calculate_text_model_cost(self, model_id: str, input_tokens: int = 0, output_tokens: int = 0, 
                                use_cached: bool = False, use_batch: bool = False, **kwargs) -> Tuple[float, Dict[str, Any]]:
        """Calculate cost for standard text-based models."""
        model = self.models[model_id]
        
        # Determine which pricing to use
        if use_batch and model.get('has_batch'):
            input_cost = model.get('batch_input_cost', model['input_cost'])
            output_cost = model.get('batch_output_cost', model['output_cost'])
            pricing_type = "Batch API"
        elif use_cached and model.get('has_cached') and model.get('cached_input_cost') is not None:
            input_cost = model['cached_input_cost']
            output_cost = model['output_cost']
            pricing_type = "Cached Input"
        else:
            input_cost = model['input_cost']
            output_cost = model['output_cost']
            pricing_type = "Standard"
        
        # Calculate costs
        input_cost_total = (input_tokens / 1_000_000) * input_cost
        output_cost_total = (output_tokens / 1_000_000) * output_cost if output_cost else 0
        total_cost = input_cost_total + output_cost_total
        
        breakdown = {
            'input_cost': input_cost_total,
            'output_cost': output_cost_total,
            'total_cost': total_cost,
            'pricing_type': pricing_type,
            'input_tokens': input_tokens,
            'output_tokens': output_tokens,
            'model_category': model.get('category', 'Unknown')
        }
        
        return total_cost, breakdown
    
    def calculate_audio_model_cost(self, model_id: str, text_input_tokens: int = 0, text_output_tokens: int = 0,
                                 audio_input_tokens: int = 0, audio_output_tokens: int = 0, **kwargs) -> Tuple[float, Dict[str, Any]]:
        """Calculate cost for audio models with separate text and audio token pricing."""
        model = self.models[model_id]
        
        # Calculate text costs
        text_input_cost = (text_input_tokens / 1_000_000) * model['text_input_cost']
        text_output_cost = (text_output_tokens / 1_000_000) * model['text_output_cost'] if model.get('text_output_cost') else 0
        
        # Calculate audio costs
        audio_input_cost = (audio_input_tokens / 1_000_000) * model['audio_input_cost'] if model.get('audio_input_cost') else 0
        audio_output_cost = (audio_output_tokens / 1_000_000) * model['audio_output_cost'] if model.get('audio_output_cost') else 0
        
        total_cost = text_input_cost + text_output_cost + audio_input_cost + audio_output_cost
        
        breakdown = {
            'text_input_cost': text_input_cost,
            'text_output_cost': text_output_cost,
            'audio_input_cost': audio_input_cost,
            'audio_output_cost': audio_output_cost,
            'total_cost': total_cost,
            'pricing_type': 'Audio Tokens',
            'text_input_tokens': text_input_tokens,
            'text_output_tokens': text_output_tokens,
            'audio_input_tokens': audio_input_tokens,
            'audio_output_tokens': audio_output_tokens
        }
        
        return total_cost, breakdown
    
    def calculate_realtime_audio_cost(self, model_id: str, text_input_tokens: int = 0, text_output_tokens: int = 0,
                                    audio_input_tokens: int = 0, audio_output_tokens: int = 0, 
                                    use_cached: bool = False, **kwargs) -> Tuple[float, Dict[str, Any]]:
        """Calculate cost for realtime audio models with caching support."""
        model = self.models[model_id]
        
        # Calculate text costs (with optional caching)
        if use_cached and model.get('text_cached_input_cost'):
            text_input_cost = (text_input_tokens / 1_000_000) * model['text_cached_input_cost']
            pricing_type = "Realtime Audio (Cached)"
        else:
            text_input_cost = (text_input_tokens / 1_000_000) * model['text_input_cost']
            pricing_type = "Realtime Audio"
            
        text_output_cost = (text_output_tokens / 1_000_000) * model['text_output_cost']
        
        # Calculate audio costs (with optional caching)
        if use_cached and model.get('audio_cached_input_cost'):
            audio_input_cost = (audio_input_tokens / 1_000_000) * model['audio_cached_input_cost']
        else:
            audio_input_cost = (audio_input_tokens / 1_000_000) * model['audio_input_cost']
            
        audio_output_cost = (audio_output_tokens / 1_000_000) * model['audio_output_cost']
        
        total_cost = text_input_cost + text_output_cost + audio_input_cost + audio_output_cost
        
        breakdown = {
            'text_input_cost': text_input_cost,
            'text_output_cost': text_output_cost,
            'audio_input_cost': audio_input_cost,
            'audio_output_cost': audio_output_cost,
            'total_cost': total_cost,
            'pricing_type': pricing_type,
            'text_input_tokens': text_input_tokens,
            'text_output_tokens': text_output_tokens,
            'audio_input_tokens': audio_input_tokens,
            'audio_output_tokens': audio_output_tokens
        }
        
        return total_cost, breakdown
    

    
    def calculate_per_minute_cost(self, model_id: str, minutes: float = 0, **kwargs) -> Tuple[float, Dict[str, Any]]:
        """Calculate cost for per-minute models (like Whisper)."""
        model = self.models[model_id]
        
        total_cost = minutes * model['cost_per_minute']
        
        breakdown = {
            'total_cost': total_cost,
            'pricing_type': 'Per Minute',
            'minutes': minutes,
            'cost_per_minute': model['cost_per_minute']
        }
        
        return total_cost, breakdown
    
    def calculate_per_character_cost(self, model_id: str, characters: int = 0, **kwargs) -> Tuple[float, Dict[str, Any]]:
        """Calculate cost for per-character models (like TTS)."""
        model = self.models[model_id]
        
        total_cost = (characters / 1_000_000) * model['cost_per_1m_characters']
        
        breakdown = {
            'total_cost': total_cost,
            'pricing_type': 'Per Character',
            'characters': characters,
            'cost_per_1m_characters': model['cost_per_1m_characters']
        }
        
        return total_cost, breakdown
    

    
    def calculate_fine_tuning_cost(self, model_id: str, training_tokens: int = 0, training_hours: float = 0,
                                 inference_input_tokens: int = 0, inference_output_tokens: int = 0,
                                 use_cached: bool = False, use_data_sharing: bool = False, **kwargs) -> Tuple[float, Dict[str, Any]]:
        """Calculate cost for fine-tuning models."""
        model = self.models[model_id]
        
        # Calculate training costs
        if 'training_cost_per_hour' in model:
            if use_data_sharing and model.get('data_sharing_training_cost_per_hour'):
                training_cost = training_hours * model['data_sharing_training_cost_per_hour']
                pricing_type = "Fine-tuning (Data Sharing)"
            else:
                training_cost = training_hours * model['training_cost_per_hour']
                pricing_type = "Fine-tuning"
        else:
            training_cost = (training_tokens / 1_000_000) * model['training_cost_per_1m_tokens']
            pricing_type = "Fine-tuning"
        
        # Calculate inference costs
        if use_data_sharing and model.get('data_sharing_inference_input_cost'):
            if use_cached and model.get('data_sharing_inference_cached_input_cost'):
                inference_input_cost = (inference_input_tokens / 1_000_000) * model['data_sharing_inference_cached_input_cost']
            else:
                inference_input_cost = (inference_input_tokens / 1_000_000) * model['data_sharing_inference_input_cost']
            inference_output_cost = (inference_output_tokens / 1_000_000) * model['data_sharing_inference_output_cost']
        else:
            if use_cached and model.get('inference_cached_input_cost'):
                inference_input_cost = (inference_input_tokens / 1_000_000) * model['inference_cached_input_cost']
            else:
                inference_input_cost = (inference_input_tokens / 1_000_000) * model['inference_input_cost']
            inference_output_cost = (inference_output_tokens / 1_000_000) * model['inference_output_cost']
        
        total_cost = training_cost + inference_input_cost + inference_output_cost
        
        breakdown = {
            'training_cost': training_cost,
            'inference_input_cost': inference_input_cost,
            'inference_output_cost': inference_output_cost,
            'total_cost': total_cost,
            'pricing_type': pricing_type,
            'training_tokens': training_tokens,
            'training_hours': training_hours,
            'inference_input_tokens': inference_input_tokens,
            'inference_output_tokens': inference_output_tokens
        }
        
        return total_cost, breakdown
    
    def calculate_free_model_cost(self, model_id: str, **kwargs) -> Tuple[float, Dict[str, Any]]:
        """Calculate cost for free models (moderation)."""
        model = self.models[model_id]
        
        breakdown = {
            'total_cost': 0.0,
            'pricing_type': 'Free',
            'model_name': model['name']
        }
        
        return 0.0, breakdown
    
    def calculate_image_generation_cost(self, model_id: str, num_images: int = 1, resolution: str = '1024x1024', 
                                      quality: str = 'standard', **kwargs) -> Tuple[float, Dict[str, Any]]:
        """Calculate cost for image generation models (per-image pricing)."""
        model = self.image_generation_models[model_id]
        
        if resolution in model['resolutions'] and quality in model['resolutions'][resolution]:
            cost_per_image = model['resolutions'][resolution][quality]
            total_cost = num_images * cost_per_image
            
            breakdown = {
                'cost_per_image': cost_per_image,
                'num_images': num_images,
                'total_cost': total_cost,
                'resolution': resolution,
                'quality': quality,
                'pricing_type': 'Per Image',
                'model_name': model['name']
            }
            
            return total_cost, breakdown
        else:
            raise ValueError(f"Invalid resolution {resolution} or quality {quality} for model {model_id}")
    
    def calculate_web_search_cost(self, model_id: str, num_calls: int = 1000, context_size: str = 'medium', **kwargs) -> Tuple[float, Dict[str, Any]]:
        """Calculate cost for web search models (per-1K-calls pricing)."""
        model = self.web_search_models[model_id]
        
        if context_size in model['search_context_pricing']:
            cost_per_1k_calls = model['search_context_pricing'][context_size]
            total_cost = (num_calls / 1000) * cost_per_1k_calls
            
            breakdown = {
                'cost_per_1k_calls': cost_per_1k_calls,
                'num_calls': num_calls,
                'total_cost': total_cost,
                'context_size': context_size,
                'pricing_type': 'Per 1K Calls',
                'model_name': model['name']
            }
            
            return total_cost, breakdown
        else:
            raise ValueError(f"Invalid context size {context_size} for model {model_id}")
    
    def calculate_tool_cost(self, tool_id: str, **kwargs) -> Tuple[float, Dict[str, Any]]:
        """Calculate cost for built-in tools."""
        tool = self.built_in_tools[tool_id]
        
        if tool_id == 'code_interpreter':
            num_containers = kwargs.get('num_containers', 1)
            total_cost = num_containers * tool['cost_per_container']
            breakdown = {
                'cost_per_container': tool['cost_per_container'],
                'num_containers': num_containers,
                'total_cost': total_cost,
                'pricing_type': 'Per Container',
                'tool_name': tool['name']
            }
        elif tool_id == 'file_search_storage':
            gb_days = kwargs.get('gb_days', 1)
            billable_gb_days = max(0, gb_days - tool['free_gb'])
            total_cost = billable_gb_days * tool['cost_per_gb_day']
            breakdown = {
                'cost_per_gb_day': tool['cost_per_gb_day'],
                'gb_days': gb_days,
                'free_gb': tool['free_gb'],
                'billable_gb_days': billable_gb_days,
                'total_cost': total_cost,
                'pricing_type': 'Per GB/Day',
                'tool_name': tool['name']
            }
        elif tool_id == 'file_search_tool_call':
            num_calls = kwargs.get('num_calls', 1)
            total_cost = (num_calls / 1000) * tool['cost_per_1k_calls']
            breakdown = {
                'cost_per_1k_calls': tool['cost_per_1k_calls'],
                'num_calls': num_calls,
                'total_cost': total_cost,
                'pricing_type': 'Per 1K Calls',
                'tool_name': tool['name']
            }
        else:
            raise ValueError(f"Unknown tool: {tool_id}")
        
        return total_cost, breakdown

    def calculate_embeddings_cost(self, model_id: str, input_tokens: int = 0, 
                                use_batch: bool = False, **kwargs) -> Tuple[float, Dict[str, Any]]:
        """Calculate cost for embeddings models."""
        model = self.models[model_id]
        
        # Use batch pricing if available and requested
        if use_batch and model.get('has_batch', False):
            cost_per_1k = model.get('batch_input_cost', model['input_cost'])
        else:
            cost_per_1k = model['input_cost']
        
        total_cost = (input_tokens / 1000) * cost_per_1k
        
        breakdown = {
            'input_cost': total_cost,
            'output_cost': 0.0,
            'total_cost': total_cost,
            'pricing_type': 'embeddings',
            'cost_per_1k_tokens': cost_per_1k,
            'dimensions': model.get('dimensions', 'N/A')
        }
        
        return total_cost, breakdown

    @staticmethod
    def get_image_models() -> Dict[str, Dict[str, Any]]:
        """Get image generation models."""
        instance = OpenAIModels()
        return instance.image_generation_models

    @staticmethod
    def get_embedding_models() -> Dict[str, Dict[str, Any]]:
        """Get embedding models from the main models collection."""
        instance = OpenAIModels()
        return {k: v for k, v in instance.models.items() if v.get('category') == 'Embeddings'}

    @staticmethod
    def validate_pricing_data() -> Dict[str, Any]:
        """Comprehensive validation of all pricing files and models."""
        validation_report = {
            'files_loaded': {},
            'model_counts': {},
            'pricing_types': {},
            'issues': [],
            'total_models': 0
        }
        
        # Expected pricing files
        expected_files = [
            'pricing_data/text_models.json',
            'pricing_data/audio_models.json',
            'pricing_data/transcription_models.json',
            'pricing_data/embeddings_models.json',
            'pricing_data/fine_tuning_models.json',
            'pricing_data/moderation_models.json',
            'pricing_data/image_generation_models.json',
            'pricing_data/web_search_models.json',
            'pricing_data/built_in_tools.json'
        ]
        
        instance = OpenAIModels()
        
        # Check each file
        for file_path in expected_files:
            file_exists = os.path.exists(file_path)
            validation_report['files_loaded'][file_path] = file_exists
            
            if file_exists:
                try:
                    with open(file_path, 'r') as f:
                        data = json.load(f)
                        for category_name, category_data in data.items():
                            count = len(category_data)
                            validation_report['model_counts'][category_name] = count
                            validation_report['total_models'] += count
                except Exception as e:
                    validation_report['issues'].append(f"Error reading {file_path}: {e}")
            else:
                validation_report['issues'].append(f"Missing file: {file_path}")
        
        # Check pricing types distribution
        pricing_type_counts = {}
        for model_id, model_info in instance.models.items():
            pricing_type = model_info.get('pricing_type', 'standard')
            pricing_type_counts[pricing_type] = pricing_type_counts.get(pricing_type, 0) + 1
        validation_report['pricing_types'] = pricing_type_counts
        
        # Check for required fields
        required_fields = ['name', 'description', 'category']
        for model_id, model_info in instance.models.items():
            for field in required_fields:
                if field not in model_info:
                    validation_report['issues'].append(f"Model {model_id} missing required field: {field}")
        
        # Check special categories
        validation_report['special_models'] = {
            'image_generation': len(instance.image_generation_models),
            'web_search': len(instance.web_search_models),
            'built_in_tools': len(instance.built_in_tools)
        }
        
        return validation_report

class FileProcessor:
    """File processing utilities for different formats."""
    
    @staticmethod
    def detect_encoding(file_bytes: bytes) -> str:
        """Detect file encoding using chardet."""
        result = chardet.detect(file_bytes)
        return result['encoding'] or 'utf-8'
    
    @staticmethod
    def extract_text_from_docx(file) -> str:
        """Extract text from DOCX files."""
        try:
            doc = docx.Document(file)
            text = []
            for paragraph in doc.paragraphs:
                text.append(paragraph.text)
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text.append(cell.text)
            
            return '\n'.join(text)
        except Exception as e:
            raise Exception(f"Error reading DOCX file: {str(e)}")
    
    @staticmethod
    def extract_text_from_pdf(file) -> str:
        """Extract text from PDF files using multiple methods."""
        text = ""
        
        try:
            # Try with pdfplumber first (better for complex layouts)
            with pdfplumber.open(file) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
        except Exception:
            # Fallback to PyPDF2
            try:
                file.seek(0)  # Reset file pointer
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
            except Exception as e:
                raise Exception(f"Error reading PDF file: {str(e)}")
        
        if not text.strip():
            raise Exception("No text could be extracted from the PDF file. It might be an image-based PDF.")
        
        return text
    
    @staticmethod
    def extract_text_from_excel(file) -> str:
        """Extract text from Excel files."""
        try:
            # Try with openpyxl first (for .xlsx files)
            try:
                workbook = openpyxl.load_workbook(file, data_only=True)
                text = []
                
                for sheet_name in workbook.sheetnames:
                    sheet = workbook[sheet_name]
                    text.append(f"Sheet: {sheet_name}")
                    
                    for row in sheet.iter_rows(values_only=True):
                        row_text = []
                        for cell in row:
                            if cell is not None:
                                row_text.append(str(cell))
                        if row_text:
                            text.append('\t'.join(row_text))
                
                return '\n'.join(text)
            
            except Exception:
                # Fallback to xlrd for older .xls files
                file.seek(0)
                workbook = xlrd.open_workbook(file_contents=file.read())
                text = []
                
                for sheet_name in workbook.sheet_names():
                    sheet = workbook.sheet_by_name(sheet_name)
                    text.append(f"Sheet: {sheet_name}")
                    
                    for row_idx in range(sheet.nrows):
                        row_text = []
                        for col_idx in range(sheet.ncols):
                            cell_value = sheet.cell_value(row_idx, col_idx)
                            if cell_value:
                                row_text.append(str(cell_value))
                        if row_text:
                            text.append('\t'.join(row_text))
                
                return '\n'.join(text)
        
        except Exception as e:
            raise Exception(f"Error reading Excel file: {str(e)}")
    
    @staticmethod
    def extract_text_from_pptx(file) -> str:
        """Extract text from PowerPoint files."""
        try:
            presentation = Presentation(file)
            text = []
            
            for i, slide in enumerate(presentation.slides, 1):
                text.append(f"Slide {i}:")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text.append(shape.text)
                
                text.append("")  # Add blank line between slides
            
            return '\n'.join(text)
        except Exception as e:
            raise Exception(f"Error reading PowerPoint file: {str(e)}")
    
    @staticmethod
    def extract_text_from_txt(file) -> str:
        """Extract text from text files with encoding detection."""
        try:
            file_bytes = file.read()
            encoding = FileProcessor.detect_encoding(file_bytes)
            return file_bytes.decode(encoding)
        except Exception as e:
            raise Exception(f"Error reading text file: {str(e)}")
    
    @staticmethod
    def process_file(uploaded_file) -> str:
        """Process uploaded file and extract text based on file type."""
        file_extension = uploaded_file.name.lower().split('.')[-1]
        
        if file_extension in ['txt', 'md', 'py', 'js', 'html', 'css', 'json', 'xml', 'csv']:
            return FileProcessor.extract_text_from_txt(uploaded_file)
        elif file_extension == 'docx':
            return FileProcessor.extract_text_from_docx(uploaded_file)
        elif file_extension == 'pdf':
            return FileProcessor.extract_text_from_pdf(uploaded_file)
        elif file_extension in ['xlsx', 'xls']:
            return FileProcessor.extract_text_from_excel(uploaded_file)
        elif file_extension == 'pptx':
            return FileProcessor.extract_text_from_pptx(uploaded_file)
        else:
            raise Exception(f"Unsupported file type: {file_extension}")

class TokenCalculator:
    """Token calculation and cost estimation utilities."""
    
    @staticmethod
    def count_tokens(text: str, model: str = "gpt-4o") -> int:
        """Count tokens in text using tiktoken."""
        try:
            # Map model names to encodings
            encoding_map = {
                # Latest models use cl100k_base
                "gpt-4": "cl100k_base",
                "gpt-4o": "cl100k_base", 
                "gpt-4o-mini": "cl100k_base",
                "gpt-3.5-turbo": "cl100k_base",
                # Reasoning models use o200k_base
                "o1": "o200k_base",
                "o1-preview": "o200k_base", 
                "o1-mini": "o200k_base",
                "o3": "o200k_base",
                "o3-mini": "o200k_base",
                "o4-mini": "o200k_base",
                # Legacy models
                "davinci": "p50k_base",
                "curie": "p50k_base",
                "babbage": "p50k_base",
                "ada": "p50k_base"
            }
            
            # Find the right encoding for the model
            encoding_name = "cl100k_base"  # default
            for model_key, encoding in encoding_map.items():
                if model_key in model.lower():
                    encoding_name = encoding
                    break
            
            # Get the encoding and count tokens
            encoding = tiktoken.get_encoding(encoding_name)
            tokens = encoding.encode(text)
            return len(tokens)
            
        except Exception as e:
            # Fallback to simple estimation: ~4 characters per token
            return len(text) // 4
    
    @staticmethod
    def calculate_cost(input_tokens: int, output_tokens: int, model_info: Dict[str, Any], 
                      use_cached: bool = False, use_batch: bool = False) -> Tuple[float, Dict[str, float]]:
        """Calculate cost for given token counts and model."""
        
        # Determine pricing based on options
        if use_batch and "batch_input_cost" in model_info:
            input_cost_per_million = model_info["batch_input_cost"]
            output_cost_per_million = model_info.get("batch_output_cost", model_info["output_cost"])
            pricing_type = "Batch API (50% discount)"
        elif use_cached and model_info.get("cached_input_cost") is not None:
            input_cost_per_million = model_info["cached_input_cost"]
            output_cost_per_million = model_info["output_cost"]
            pricing_type = "Cached input (50% discount)"
        else:
            input_cost_per_million = model_info["input_cost"]
            output_cost_per_million = model_info["output_cost"]
            pricing_type = "Standard pricing"
        
        # Calculate costs
        input_cost = (input_tokens / 1_000_000) * input_cost_per_million
        output_cost = (output_tokens / 1_000_000) * output_cost_per_million if output_cost_per_million else 0
        total_cost = input_cost + output_cost
        
        breakdown = {
            "input_cost": input_cost,
            "output_cost": output_cost,
            "total_cost": total_cost,
            "pricing_type": pricing_type
        }
        
        return total_cost, breakdown

def create_cost_comparison_chart(comparisons: List[Dict[str, Any]]) -> go.Figure:
    """Create a cost comparison chart for multiple models."""
    if not comparisons:
        return go.Figure()
    
    models = [comp['model'] for comp in comparisons]
    costs = [comp['cost'] for comp in comparisons]
    
    fig = go.Figure(data=[
        go.Bar(x=models, y=costs, text=[f"${cost:.4f}" for cost in costs], textposition='auto')
    ])
    
    fig.update_layout(
        title="Cost Comparison Across Models",
        xaxis_title="Models",
        yaxis_title="Cost (USD)",
        showlegend=False
    )
    
    return fig

def create_token_breakdown_chart(input_tokens: int, output_tokens: int) -> go.Figure:
    """Create a pie chart showing token breakdown."""
    labels = ['Input Tokens', 'Output Tokens']
    values = [input_tokens, output_tokens]
    
    fig = go.Figure(data=[go.Pie(labels=labels, values=values)])
    fig.update_layout(title="Token Distribution")
    
    return fig

def main():
    """Main Streamlit application."""
    
    # Initialize the models instance
    openai_models = OpenAIModels()
    
    # Header
    st.markdown('<h1 class="main-header">🧮 OpenAI Token Counter & Cost Calculator</h1>', unsafe_allow_html=True)
    st.markdown("Calculate token counts and costs for all OpenAI models with support for different pricing categories.")
    
    # Sidebar configuration
    st.sidebar.header("⚙️ Configuration")
    
    # Model selection with categories
    all_models = openai_models.models
    categories = list(set(model.get('category', 'Unknown') for model in all_models.values()))
    categories.sort()
    
    selected_category = st.sidebar.selectbox("Select Model Category", categories)
    
    # Filter models by category
    category_models = {k: v for k, v in all_models.items() if v.get('category') == selected_category}
    model_options = list(category_models.keys())
    
    selected_model = st.sidebar.selectbox("Select Model", model_options)
    
    if selected_model and selected_model in category_models:
        model_info = category_models[selected_model]
        
        # Display model information
        with st.sidebar.expander("📊 Model Information", expanded=True):
            st.write(f"**Name:** {model_info['name']}")
            st.write(f"**Description:** {model_info['description']}")
            st.write(f"**Category:** {model_info.get('category', 'Unknown')}")
            
            if 'context_window' in model_info:
                st.write(f"**Context Window:** {model_info['context_window']:,} tokens")
            
            # Pricing information - handle different pricing structures
            st.write("**Pricing:**")
            
            # Standard text model pricing
            if 'input_cost' in model_info:
                st.write(f"- Input: ${model_info['input_cost']:.2f}/1M tokens")
                if model_info.get('output_cost'):
                    st.write(f"- Output: ${model_info['output_cost']:.2f}/1M tokens")
                if model_info.get('cached_input_cost'):
                    st.write(f"- Cached Input: ${model_info['cached_input_cost']:.2f}/1M tokens")
            
            # Audio model pricing
            elif 'text_input_cost' in model_info:
                st.write(f"- Text Input: ${model_info['text_input_cost']:.2f}/1M tokens")
                if model_info.get('text_output_cost'):
                    st.write(f"- Text Output: ${model_info['text_output_cost']:.2f}/1M tokens")
                if model_info.get('audio_input_cost'):
                    st.write(f"- Audio Input: ${model_info['audio_input_cost']:.2f}/1M tokens")
                if model_info.get('audio_output_cost'):
                    st.write(f"- Audio Output: ${model_info['audio_output_cost']:.2f}/1M tokens")
            
            # Fine-tuning pricing
            elif 'training_cost_per_1m_tokens' in model_info:
                st.write(f"- Training: ${model_info['training_cost_per_1m_tokens']:.2f}/1M tokens")
                if model_info.get('inference_input_cost'):
                    st.write(f"- Inference Input: ${model_info['inference_input_cost']:.2f}/1M tokens")
                if model_info.get('inference_output_cost'):
                    st.write(f"- Inference Output: ${model_info['inference_output_cost']:.2f}/1M tokens")
            
            # Training cost per hour (for some models)
            elif 'training_cost_per_hour' in model_info:
                st.write(f"- Training: ${model_info['training_cost_per_hour']:.2f}/hour")
                if model_info.get('inference_input_cost'):
                    st.write(f"- Inference Input: ${model_info['inference_input_cost']:.2f}/1M tokens")
            
            # Per minute costs
            elif 'cost_per_minute' in model_info:
                st.write(f"- Cost: ${model_info['cost_per_minute']:.3f}/minute")
            
            # Per character costs
            elif 'cost_per_1m_characters' in model_info:
                st.write(f"- Cost: ${model_info['cost_per_1m_characters']:.2f}/1M characters")
            
            # Free models
            elif model_info.get('is_free'):
                st.write("- **FREE** ✨")
            
            else:
                st.write("- See special pricing in Tools tab")
    
    # Pricing options
    st.sidebar.subheader("💰 Pricing Options")
    use_cached = st.sidebar.checkbox("Use Cached Input (50% discount)", 
                                   disabled=not (selected_model and category_models.get(selected_model, {}).get('cached_input_cost')))
    use_batch = st.sidebar.checkbox("Use Batch API (50% discount)", 
                                  disabled=not (selected_model and category_models.get(selected_model, {}).get('has_batch')))
    
    # Main content area with tabs
    tab1, tab2, tab3, tab4, tab5, tab6, tab7 = st.tabs(["📝 Text Input", "📁 File Upload", "📊 Analytics", "🔧 Tools & Special Models", "🔍 System Status", "🧪 Model Testing Hub", "⚙️ Model Management"])
    
    with tab1:
        st.subheader("Enter Text for Token Counting")
        
        if selected_model and selected_model in category_models:
            model_info = category_models[selected_model]
            pricing_type = model_info.get('pricing_type', 'standard')
            
            # Input form based on pricing type
            if pricing_type == 'audio_tokens' or pricing_type == 'realtime_audio':
                # Audio models - separate text and audio inputs
                input_text = st.text_area("Input Text", height=150, placeholder="Enter your text here...")
                
                col1, col2 = st.columns(2)
                with col1:
                    text_output_tokens = st.number_input("Text Output Tokens", min_value=0, value=100, step=10)
                    audio_input_tokens = st.number_input("Audio Input Tokens", min_value=0, value=0, step=100)
                with col2:
                    audio_output_tokens = st.number_input("Audio Output Tokens", min_value=0, value=0, step=100)
                
                if input_text:
                    text_input_tokens = TokenCalculator.count_tokens(input_text, selected_model)
                    total_cost, breakdown = openai_models.calculate_cost_by_category(
                        selected_model,
                        text_input_tokens=text_input_tokens,
                        text_output_tokens=text_output_tokens,
                        audio_input_tokens=audio_input_tokens,
                        audio_output_tokens=audio_output_tokens,
                        use_cached=use_cached
                    )
                    
                    col1, col2, col3, col4 = st.columns(4)
                    with col1:
                        st.metric("Text Input", f"{text_input_tokens:,}")
                    with col2:
                        st.metric("Text Output", f"{text_output_tokens:,}")
                    with col3:
                        st.metric("Audio Tokens", f"{audio_input_tokens + audio_output_tokens:,}")
                    with col4:
                        st.metric("Total Cost", f"${total_cost:.6f}")
                    
                    # Cost breakdown
                    with st.expander("💰 Cost Breakdown", expanded=True):
                        for key, value in breakdown.items():
                            if key.endswith('_cost') and value > 0:
                                st.write(f"**{key.replace('_', ' ').title()}:** ${value:.6f}")
                            elif key == 'pricing_type':
                                st.write(f"**{key.replace('_', ' ').title()}:** {value}")
            
            elif pricing_type == 'transcription' or pricing_type == 'tts':
                # Transcription/TTS models - tokens + minutes
                input_text = st.text_area("Input Text", height=150, placeholder="Enter your text here...")
                
                col1, col2 = st.columns(2)
                with col1:
                    minutes = st.number_input("Audio Duration (minutes)", min_value=0.0, value=1.0, step=0.1)
                with col2:
                    if pricing_type == 'transcription':
                        audio_input_tokens = st.number_input("Audio Input Tokens", min_value=0, value=0, step=100)
                        text_output_tokens = st.number_input("Text Output Tokens", min_value=0, value=100, step=10)
                    else:
                        audio_output_tokens = st.number_input("Audio Output Tokens", min_value=0, value=0, step=100)
                
                if input_text:
                    text_input_tokens = TokenCalculator.count_tokens(input_text, selected_model)
                    kwargs = {
                        'text_input_tokens': text_input_tokens,
                        'minutes': minutes
                    }
                    if pricing_type == 'transcription':
                        kwargs.update({
                            'audio_input_tokens': audio_input_tokens,
                            'text_output_tokens': text_output_tokens
                        })
                    else:
                        kwargs.update({
                            'audio_output_tokens': audio_output_tokens
                        })
                    
                    total_cost, breakdown = openai_models.calculate_cost_by_category(selected_model, **kwargs)
                    
                    col1, col2, col3 = st.columns(3)
                    with col1:
                        st.metric("Text Tokens", f"{text_input_tokens:,}")
                    with col2:
                        st.metric("Duration", f"{minutes:.1f} min")
                    with col3:
                        st.metric("Total Cost", f"${total_cost:.6f}")
            
            elif pricing_type == 'per_minute':
                # Per-minute pricing (like Whisper)
                col1, col2 = st.columns(2)
                with col1:
                    minutes = st.number_input("Audio Duration (minutes)", min_value=0.0, value=1.0, step=0.1)
                
                total_cost, breakdown = openai_models.calculate_cost_by_category(selected_model, minutes=minutes)
                
                col1, col2 = st.columns(2)
                with col1:
                    st.metric("Duration", f"{minutes:.1f} min")
                with col2:
                    st.metric("Total Cost", f"${total_cost:.6f}")
            
            elif pricing_type == 'per_character':
                # Per-character pricing (like TTS)
                input_text = st.text_area("Input Text", height=150, placeholder="Enter your text here...")
                
                if input_text:
                    characters = len(input_text)
                    total_cost, breakdown = openai_models.calculate_cost_by_category(selected_model, characters=characters)
                    
                    col1, col2 = st.columns(2)
                    with col1:
                        st.metric("Characters", f"{characters:,}")
                    with col2:
                        st.metric("Total Cost", f"${total_cost:.6f}")
            
            elif pricing_type == 'embeddings':
                # Embeddings models
                input_text = st.text_area("Input Text", height=150, placeholder="Enter your text here...")
                
                if input_text:
                    input_tokens = TokenCalculator.count_tokens(input_text, selected_model)
                    total_cost, breakdown = openai_models.calculate_cost_by_category(
                        selected_model, 
                        input_tokens=input_tokens, 
                        use_batch=use_batch
                    )
                    
                    col1, col2, col3 = st.columns(3)
                    with col1:
                        st.metric("Input Tokens", f"{input_tokens:,}")
                    with col2:
                        st.metric("Dimensions", breakdown.get('dimensions', 'N/A'))
                    with col3:
                        st.metric("Total Cost", f"${total_cost:.6f}")
            
            elif pricing_type == 'fine_tuning':
                # Fine-tuning models
                st.info("Fine-tuning pricing includes both training and inference costs.")
                
                col1, col2 = st.columns(2)
                with col1:
                    if 'training_cost_per_hour' in model_info:
                        training_hours = st.number_input("Training Hours", min_value=0.0, value=1.0, step=0.1)
                        training_tokens = 0
                    else:
                        training_tokens = st.number_input("Training Tokens", min_value=0, value=1000000, step=100000)
                        training_hours = 0
                
                with col2:
                    inference_input_tokens = st.number_input("Inference Input Tokens", min_value=0, value=1000, step=100)
                    inference_output_tokens = st.number_input("Inference Output Tokens", min_value=0, value=100, step=10)
                
                use_data_sharing = st.checkbox("Use Data Sharing (reduced costs)", 
                                             disabled=not model_info.get('has_data_sharing'))
                
                total_cost, breakdown = openai_models.calculate_cost_by_category(
                    selected_model,
                    training_tokens=training_tokens,
                    training_hours=training_hours,
                    inference_input_tokens=inference_input_tokens,
                    inference_output_tokens=inference_output_tokens,
                    use_cached=use_cached,
                    use_data_sharing=use_data_sharing
                )
                
                col1, col2, col3 = st.columns(3)
                with col1:
                    if training_hours > 0:
                        st.metric("Training Hours", f"{training_hours:.1f}")
                    else:
                        st.metric("Training Tokens", f"{training_tokens:,}")
                with col2:
                    st.metric("Inference Tokens", f"{inference_input_tokens + inference_output_tokens:,}")
                with col3:
                    st.metric("Total Cost", f"${total_cost:.6f}")
            
            elif model_info.get('is_free'):
                # Free models
                input_text = st.text_area("Input Text", height=150, placeholder="Enter your text here...")
                
                if input_text:
                    tokens = TokenCalculator.count_tokens(input_text, selected_model)
                    total_cost, breakdown = openai_models.calculate_cost_by_category(selected_model)
                    
                    col1, col2 = st.columns(2)
                    with col1:
                        st.metric("Tokens", f"{tokens:,}")
                    with col2:
                        st.metric("Total Cost", "FREE ✨")
            
            else:
                # Standard text models
                input_text = st.text_area("Input Text", height=150, placeholder="Enter your text here...")
                
                col1, col2 = st.columns(2)
                with col1:
                    estimated_output_tokens = st.number_input("Estimated Output Tokens", min_value=0, value=100, step=10)
                
                if input_text:
                    input_tokens = TokenCalculator.count_tokens(input_text, selected_model)
                    total_cost, breakdown = openai_models.calculate_cost_by_category(
                        selected_model,
                        input_tokens=input_tokens,
                        output_tokens=estimated_output_tokens,
                        use_cached=use_cached,
                        use_batch=use_batch
                    )
                    
                    col1, col2, col3 = st.columns(3)
                    with col1:
                        st.metric("Input Tokens", f"{input_tokens:,}")
                    with col2:
                        st.metric("Output Tokens", f"{estimated_output_tokens:,}")
                    with col3:
                        st.metric("Total Cost", f"${total_cost:.6f}")
                    
                    # Context window warning
                    if 'context_window' in model_info and (input_tokens + estimated_output_tokens) > model_info['context_window']:
                        st.warning(f"⚠️ Token count ({input_tokens + estimated_output_tokens:,}) exceeds model's context window ({model_info['context_window']:,})")
            
            # Display cost breakdown for all pricing types
            if 'breakdown' in locals() and breakdown:
                with st.expander("💰 Cost Breakdown", expanded=True):
                    for key, value in breakdown.items():
                        if key == 'pricing_type':
                            st.write(f"**{key.replace('_', ' ').title()}:** {value}")
                        elif key.endswith('_cost') and isinstance(value, (int, float)):
                            st.write(f"**{key.replace('_', ' ').title()}:** ${value:.6f}")
                        elif key in ['minutes', 'characters', 'num_images', 'num_calls'] and isinstance(value, (int, float)):
                            st.write(f"**{key.replace('_', ' ').title()}:** {value:,}")
        else:
            st.info("Please select a model to calculate costs.")
    
    with tab2:
        st.subheader("Upload Files for Processing")
        
        uploaded_files = st.file_uploader(
            "Choose files",
            type=['txt', 'pdf', 'docx', 'xlsx', 'xls', 'pptx', 'md', 'py', 'js', 'html', 'css', 'json', 'xml', 'csv'],
            accept_multiple_files=True
        )
        
        if uploaded_files:
            for uploaded_file in uploaded_files:
                with st.expander(f"📄 {uploaded_file.name}", expanded=len(uploaded_files) == 1):
                    try:
                        # Process file
                        with st.spinner(f"Processing {uploaded_file.name}..."):
                            extracted_text = FileProcessor.process_file(uploaded_file)
                        
                        # Show preview
                        st.text_area("Extracted Text Preview", extracted_text[:500] + "..." if len(extracted_text) > 500 else extracted_text, height=150)
                        
                        if selected_model:
                            # Calculate tokens and cost
                            input_tokens = TokenCalculator.count_tokens(extracted_text, selected_model)
                            estimated_output = st.number_input(f"Estimated Output Tokens for {uploaded_file.name}", min_value=0, value=100, step=10, key=f"output_{uploaded_file.name}")
                            
                            try:
                                total_cost, breakdown = openai_models.calculate_cost_by_category(
                                    selected_model,
                                    input_tokens=input_tokens,
                                    output_tokens=estimated_output,
                                    use_cached=use_cached,
                                    use_batch=use_batch
                                )
                                
                                # Display results
                                col1, col2, col3 = st.columns(3)
                                with col1:
                                    st.metric("Input Tokens", f"{input_tokens:,}")
                                with col2:
                                    st.metric("Output Tokens", f"{estimated_output:,}")
                                with col3:
                                    st.metric("Total Cost", f"${total_cost:.6f}")
                                
                                # Show pricing type
                                st.caption(f"Pricing: {breakdown.get('pricing_type', 'Standard')}")
                            except Exception as e:
                                st.error(f"Error calculating cost for {selected_model}: {e}")
                                # Fallback to basic calculation
                                st.metric("Input Tokens", f"{input_tokens:,}")
                                st.metric("Output Tokens", f"{estimated_output:,}")
                    
                    except Exception as e:
                        st.error(f"Error processing {uploaded_file.name}: {str(e)}")
    
    with tab3:
        st.subheader("📊 Cost Analytics")
        
        # Model comparison
        st.write("### Model Cost Comparison")
        comparison_text = st.text_area("Text for Comparison", value="Hello, how are you today?", height=100)
        comparison_output_tokens = st.number_input("Output Tokens for Comparison", min_value=0, value=50, step=10)
        
        if comparison_text:
            comparison_results = []
            input_tokens = TokenCalculator.count_tokens(comparison_text, "gpt-4o")  # Use standard model for tokenization
            
            # Compare top models
            top_models = ["gpt-4o-mini-2024-07-18", "gpt-4o-2024-11-20", "o1-mini-2024-09-12"]
            available_models = [m for m in top_models if m in openai_models.models]
            
            for model_id in available_models:
                try:
                    cost, _ = openai_models.calculate_cost_by_category(
                        model_id, 
                        input_tokens=input_tokens, 
                        output_tokens=comparison_output_tokens
                    )
                    comparison_results.append({
                        'model': openai_models.models[model_id]['name'],
                        'cost': cost
                    })
                except Exception as e:
                    # Skip models that can't be calculated
                    continue
            
            if comparison_results:
                # Create comparison chart
                fig = create_cost_comparison_chart(comparison_results)
                st.plotly_chart(fig, use_container_width=True)
                
                # Show comparison table
                df = pd.DataFrame(comparison_results)
                st.dataframe(df, use_container_width=True)
    
    with tab4:
        st.subheader("🔧 Special Models & Tools")
        
        # Create three columns for different categories
        col1, col2, col3 = st.columns(3)
        
        with col1:
            # Image Generation Models
            if openai_models.image_generation_models:
                st.write("### 🎨 Image Generation")
                
                img_model = st.selectbox("Image Model", list(openai_models.image_generation_models.keys()), key="img_model")
                if img_model:
                    model_data = openai_models.image_generation_models[img_model]
                    
                    num_images = st.number_input("Number of Images", min_value=1, value=1, key="num_images")
                    resolution = st.selectbox("Resolution", list(model_data['resolutions'].keys()), key="resolution")
                    quality = st.selectbox("Quality", list(model_data['resolutions'][resolution].keys()), key="quality")
                    
                    try:
                        cost, breakdown = openai_models.calculate_image_generation_cost(
                            img_model, num_images=num_images, resolution=resolution, quality=quality
                        )
                        st.metric("Total Cost", f"${cost:.4f}")
                        with st.expander("Image Cost Details"):
                            st.write(f"**Cost per image:** ${breakdown['cost_per_image']:.4f}")
                            st.write(f"**Resolution:** {breakdown['resolution']}")
                            st.write(f"**Quality:** {breakdown['quality']}")
                    except Exception as e:
                        st.error(f"Error calculating image cost: {e}")
        
        with col2:
            # Web Search Models
            if openai_models.web_search_models:
                st.write("### 🔍 Web Search")
                
                search_model = st.selectbox("Search Model", list(openai_models.web_search_models.keys()), key="search_model")
                if search_model:
                    model_data = openai_models.web_search_models[search_model]
                    
                    num_calls = st.number_input("Number of Calls", min_value=1, value=1000, key="num_calls")
                    context_size = st.selectbox("Context Size", list(model_data['search_context_pricing'].keys()), key="context_size")
                    
                    try:
                        cost, breakdown = openai_models.calculate_web_search_cost(
                            search_model, num_calls=num_calls, context_size=context_size
                        )
                        st.metric("Total Cost", f"${cost:.4f}")
                        with st.expander("Search Cost Details"):
                            st.write(f"**Cost per 1K calls:** ${breakdown['cost_per_1k_calls']:.2f}")
                            st.write(f"**Number of calls:** {breakdown['num_calls']:,}")
                            st.write(f"**Context size:** {breakdown['context_size']}")
                    except Exception as e:
                        st.error(f"Error calculating search cost: {e}")
        
        with col3:
            # Built-in Tools
            if openai_models.built_in_tools:
                st.write("### 🛠️ Built-in Tools")
                
                tool_id = st.selectbox("Tool", list(openai_models.built_in_tools.keys()), key="tool_id")
                if tool_id:
                    tool_data = openai_models.built_in_tools[tool_id]
                    
                    if tool_id == 'code_interpreter':
                        num_containers = st.number_input("Number of Containers", min_value=1, value=1, key="containers")
                        kwargs = {'num_containers': num_containers}
                    elif tool_id == 'file_search_storage':
                        gb_days = st.number_input("GB-Days", min_value=0.0, value=1.0, step=0.1, key="gb_days")
                        kwargs = {'gb_days': gb_days}
                    elif tool_id == 'file_search_tool_call':
                        num_calls = st.number_input("Number of Calls", min_value=1, value=1000, key="tool_calls")
                        kwargs = {'num_calls': num_calls}
                    
                    try:
                        cost, breakdown = openai_models.calculate_tool_cost(tool_id, **kwargs)
                        st.metric("Total Cost", f"${cost:.4f}")
                        with st.expander("Tool Cost Details"):
                            for key, value in breakdown.items():
                                if key != 'total_cost' and key != 'tool_name':
                                    if isinstance(value, (int, float)):
                                        if 'cost' in key:
                                            st.write(f"**{key.replace('_', ' ').title()}:** ${value:.4f}")
                                        else:
                                            st.write(f"**{key.replace('_', ' ').title()}:** {value:,}")
                                    else:
                                        st.write(f"**{key.replace('_', ' ').title()}:** {value}")
                    except Exception as e:
                        st.error(f"Error calculating tool cost: {e}")
        
        # Comprehensive pricing summary
        st.write("---")
        st.write("### 📋 Comprehensive Model Summary")
        
        # Create a summary table of all available models by category
        summary_data = []
        
        # Regular models
        for model_id, model_info in openai_models.models.items():
            category = model_info.get('category', 'Unknown')
            pricing_type = model_info.get('pricing_type', 'standard')
            
            price_info = ""
            if model_info.get('is_free'):
                price_info = "FREE"
            elif 'input_cost' in model_info:
                price_info = f"${model_info['input_cost']:.2f}/1M input"
            elif 'text_input_cost' in model_info:
                price_info = f"${model_info['text_input_cost']:.2f}/1M text"
            elif 'cost_per_minute' in model_info:
                price_info = f"${model_info['cost_per_minute']:.3f}/min"
            elif 'cost_per_1m_characters' in model_info:
                price_info = f"${model_info['cost_per_1m_characters']:.2f}/1M chars"
            elif 'training_cost_per_1m_tokens' in model_info:
                price_info = f"${model_info['training_cost_per_1m_tokens']:.2f}/1M train"
            elif 'training_cost_per_hour' in model_info:
                price_info = f"${model_info['training_cost_per_hour']:.0f}/hour"
            
            summary_data.append({
                'Model': model_info['name'],
                'Category': category,
                'Pricing Type': pricing_type.title(),
                'Base Price': price_info,
                'Context Window': f"{model_info.get('context_window', 'N/A'):,}" if model_info.get('context_window') else 'N/A'
            })
        
        # Image generation models
        for model_id, model_info in openai_models.image_generation_models.items():
            summary_data.append({
                'Model': model_info['name'],
                'Category': 'Image Generation',
                'Pricing Type': 'Per Image',
                'Base Price': 'Varies by resolution',
                'Context Window': 'N/A'
            })
        
        # Web search models
        for model_id, model_info in openai_models.web_search_models.items():
            summary_data.append({
                'Model': model_info['name'],
                'Category': 'Web Search',
                'Pricing Type': 'Per 1K Calls',
                'Base Price': 'Varies by context',
                'Context Window': f"{model_info.get('context_window', 'N/A'):,}" if model_info.get('context_window') else 'N/A'
            })
        
        # Built-in tools
        for tool_id, tool_info in openai_models.built_in_tools.items():
            price_info = ""
            if 'cost_per_container' in tool_info:
                price_info = f"${tool_info['cost_per_container']:.2f}/container"
            elif 'cost_per_gb_day' in tool_info:
                price_info = f"${tool_info['cost_per_gb_day']:.3f}/GB-day"
            elif 'cost_per_1k_calls' in tool_info:
                price_info = f"${tool_info['cost_per_1k_calls']:.2f}/1K calls"
            
            summary_data.append({
                'Model': tool_info['name'],
                'Category': 'Built-in Tools',
                'Pricing Type': tool_info.get('pricing_type', 'Special').title(),
                'Base Price': price_info,
                'Context Window': 'N/A'
            })
        
        # Display summary table
        if summary_data:
            df = pd.DataFrame(summary_data)
            st.dataframe(df, use_container_width=True)
    
    with tab5:
        st.subheader("🔍 System Status & Validation")
        
        # Run comprehensive validation
        with st.spinner("Running comprehensive system validation..."):
            validation_report = openai_models.validate_pricing_data()
        
        # Status overview
        col1, col2, col3, col4 = st.columns(4)
        
        with col1:
            files_loaded = sum(1 for v in validation_report['files_loaded'].values() if v)
            total_files = len(validation_report['files_loaded'])
            st.metric("Pricing Files Loaded", f"{files_loaded}/{total_files}")
        
        with col2:
            st.metric("Total Models", validation_report['total_models'])
        
        with col3:
            issues_count = len(validation_report['issues'])
            st.metric("Issues Found", issues_count)
        
        with col4:
            special_count = sum(validation_report['special_models'].values())
            st.metric("Special Models", special_count)
        
        # Detailed file status
        st.write("### 📄 Pricing Files Status")
        file_status_data = []
        for file_path, loaded in validation_report['files_loaded'].items():
            status = "✅ Loaded" if loaded else "❌ Missing"
            file_status_data.append({
                'File': file_path.replace('pricing_data/', ''),
                'Status': status
            })
        
        file_df = pd.DataFrame(file_status_data)
        st.dataframe(file_df, use_container_width=True)
        
        # Model counts by category
        if validation_report['model_counts']:
            st.write("### 📊 Models by Category")
            category_data = []
            for category, count in validation_report['model_counts'].items():
                category_data.append({
                    'Category': category.replace('_', ' ').title(),
                    'Model Count': count
                })
            
            category_df = pd.DataFrame(category_data)
            
            # Create a bar chart
            fig = px.bar(category_df, x='Category', y='Model Count', 
                        title="Model Distribution by Category",
                        color='Model Count',
                        color_continuous_scale='viridis')
            fig.update_layout(xaxis_tickangle=-45)
            st.plotly_chart(fig, use_container_width=True)
            
            st.dataframe(category_df, use_container_width=True)
        
        # Pricing types distribution
        if validation_report['pricing_types']:
            st.write("### 💰 Pricing Types Distribution")
            pricing_data = []
            for pricing_type, count in validation_report['pricing_types'].items():
                pricing_data.append({
                    'Pricing Type': pricing_type.replace('_', ' ').title(),
                    'Count': count
                })
            
            pricing_df = pd.DataFrame(pricing_data)
            
            # Create a pie chart
            fig = px.pie(pricing_df, values='Count', names='Pricing Type',
                        title="Distribution of Pricing Types")
            st.plotly_chart(fig, use_container_width=True)
            
            st.dataframe(pricing_df, use_container_width=True)
        
        # Issues and warnings
        if validation_report['issues']:
            st.write("### ⚠️ Issues Found")
            for issue in validation_report['issues']:
                st.error(issue)
        else:
            st.success("✅ No issues found! All pricing files are properly loaded and validated.")
        
        # System information
        st.write("### 🖥️ System Information")
        system_info = {
            'Python Libraries': {
                'Streamlit': st.__version__,
                'Pandas': pd.__version__,
                'NumPy': np.__version__,
                'Tiktoken': tiktoken.__version__ if hasattr(tiktoken, '__version__') else 'Available'
            },
            'File Support': {
                'PDF Processing': 'PyPDF2 + pdfplumber',
                'Office Documents': 'python-docx, openpyxl, python-pptx',
                'Text Encoding': 'chardet',
                'Image Processing': 'Pillow'
            }
        }
        
        for category, info in system_info.items():
            with st.expander(f"📋 {category}"):
                for key, value in info.items():
                    st.write(f"**{key}:** {value}")
    
    with tab6:
        st.subheader("🧪 Comprehensive Model Testing Hub")
        st.markdown("Test and calculate costs for all OpenAI model types with dedicated interfaces.")
        
        # Create testing categories
        testing_category = st.selectbox(
            "Choose Testing Category",
            ["🤖 Text & Chat Models", "🎵 Audio & Speech Models", "🔤 Embeddings Models", 
             "🎨 Image Generation", "🔍 Search & Tools", "🎓 Fine-tuning Models", "🆓 Free Models"],
            key="testing_category"
        )
        
        if testing_category == "🤖 Text & Chat Models":
            st.write("### Text and Chat Models Testing")
            
            # Filter text models
            text_models = {k: v for k, v in openai_models.models.items() 
                          if v.get('pricing_type', 'standard') == 'standard'}
            
            if text_models:
                model_id = st.selectbox("Select Text Model", list(text_models.keys()), key="text_test_model")
                model_info = text_models[model_id]
                
                col1, col2 = st.columns([2, 1])
                with col1:
                    test_input = st.text_area("Input Text", height=200, 
                                            placeholder="Enter your text here to test token counting and cost calculation...",
                                            key="text_test_input")
                with col2:
                    output_tokens = st.number_input("Expected Output Tokens", min_value=0, value=100, step=10, key="text_output")
                    use_cache = st.checkbox("Use Cached Input", key="text_cache")
                    use_batch_api = st.checkbox("Use Batch API", key="text_batch")
                
                if test_input:
                    input_tokens = TokenCalculator.count_tokens(test_input, model_id)
                    cost, breakdown = openai_models.calculate_text_model_cost(
                        model_id, input_tokens, output_tokens, use_cache, use_batch_api
                    )
                    
                    # Results
                    col1, col2, col3, col4 = st.columns(4)
                    with col1:
                        st.metric("Input Tokens", f"{input_tokens:,}")
                    with col2:
                        st.metric("Output Tokens", f"{output_tokens:,}")
                    with col3:
                        st.metric("Total Tokens", f"{input_tokens + output_tokens:,}")
                    with col4:
                        st.metric("Total Cost", f"${cost:.6f}")
                    
                    # Cost breakdown
                    with st.expander("💰 Detailed Cost Breakdown"):
                        for key, value in breakdown.items():
                            if isinstance(value, (int, float)) and 'cost' in key:
                                st.write(f"**{key.replace('_', ' ').title()}:** ${value:.6f}")
                            elif key == 'pricing_type':
                                st.write(f"**{key.replace('_', ' ').title()}:** {value}")
        
        elif testing_category == "🎵 Audio & Speech Models":
            st.write("### Audio and Speech Models Testing")
            
            # Audio models subcategory
            audio_subcategory = st.radio("Audio Model Type", 
                                       ["🎙️ Audio Processing", "🗣️ Text-to-Speech", "📝 Speech-to-Text"],
                                       key="audio_subcategory")
            
            if audio_subcategory == "🎙️ Audio Processing":
                # Filter audio models
                audio_models = {k: v for k, v in openai_models.models.items() 
                               if v.get('pricing_type') in ['audio_tokens', 'realtime_audio']}
                
                if audio_models:
                    model_id = st.selectbox("Select Audio Model", list(audio_models.keys()), key="audio_test_model")
                    
                    col1, col2 = st.columns(2)
                    with col1:
                        st.write("**Text Components:**")
                        text_input = st.text_area("Text Input", height=100, key="audio_text_input")
                        text_output_tokens = st.number_input("Text Output Tokens", min_value=0, value=100, key="audio_text_output")
                    
                    with col2:
                        st.write("**Audio Components:**")
                        audio_input_tokens = st.number_input("Audio Input Tokens", min_value=0, value=1000, step=100, key="audio_input")
                        audio_output_tokens = st.number_input("Audio Output Tokens", min_value=0, value=500, step=100, key="audio_output")
                    
                    use_cached = st.checkbox("Use Cached Input (for supported models)", key="audio_cached")
                    
                    if text_input or audio_input_tokens > 0:
                        text_input_tokens = TokenCalculator.count_tokens(text_input, model_id) if text_input else 0
                        
                        cost, breakdown = openai_models.calculate_cost_by_category(
                            model_id,
                            text_input_tokens=text_input_tokens,
                            text_output_tokens=text_output_tokens,
                            audio_input_tokens=audio_input_tokens,
                            audio_output_tokens=audio_output_tokens,
                            use_cached=use_cached
                        )
                        
                        # Results
                        col1, col2, col3 = st.columns(3)
                        with col1:
                            st.metric("Text Tokens", f"{text_input_tokens + text_output_tokens:,}")
                        with col2:
                            st.metric("Audio Tokens", f"{audio_input_tokens + audio_output_tokens:,}")
                        with col3:
                            st.metric("Total Cost", f"${cost:.6f}")
            
            elif audio_subcategory == "🗣️ Text-to-Speech":
                # TTS models
                tts_models = {k: v for k, v in openai_models.models.items() 
                             if v.get('pricing_type') == 'per_character'}
                
                if tts_models:
                    model_id = st.selectbox("Select TTS Model", list(tts_models.keys()), key="tts_test_model")
                    
                    tts_text = st.text_area("Text to Convert to Speech", height=150, 
                                          placeholder="Enter text you want to convert to speech...",
                                          key="tts_text_input")
                    
                    if tts_text:
                        characters = len(tts_text)
                        cost, breakdown = openai_models.calculate_per_character_cost(model_id, characters)
                        
                        col1, col2, col3 = st.columns(3)
                        with col1:
                            st.metric("Characters", f"{characters:,}")
                        with col2:
                            st.metric("Cost/1M chars", f"${breakdown['cost_per_1m_characters']:.2f}")
                        with col3:
                            st.metric("Total Cost", f"${cost:.6f}")
            
            elif audio_subcategory == "📝 Speech-to-Text":
                # Whisper models
                whisper_models = {k: v for k, v in openai_models.models.items() 
                                if v.get('pricing_type') == 'per_minute'}
                
                if whisper_models:
                    model_id = st.selectbox("Select Transcription Model", list(whisper_models.keys()), key="whisper_test_model")
                    
                    col1, col2 = st.columns(2)
                    with col1:
                        audio_duration = st.number_input("Audio Duration (minutes)", min_value=0.0, value=1.0, step=0.1, key="whisper_duration")
                        st.info("Whisper models charge per minute of audio processed")
                    
                    with col2:
                        if audio_duration > 0:
                            cost, breakdown = openai_models.calculate_per_minute_cost(model_id, audio_duration)
                            
                            st.metric("Duration", f"{audio_duration:.1f} min")
                            st.metric("Cost/minute", f"${breakdown['cost_per_minute']:.3f}")
                            st.metric("Total Cost", f"${cost:.6f}")
        
        elif testing_category == "🔤 Embeddings Models":
            st.write("### Embeddings Models Testing")
            
            embeddings_models = {k: v for k, v in openai_models.models.items() 
                               if v.get('pricing_type') == 'embeddings'}
            
            if embeddings_models:
                model_id = st.selectbox("Select Embeddings Model", list(embeddings_models.keys()), key="emb_test_model")
                
                col1, col2 = st.columns([2, 1])
                with col1:
                    embed_text = st.text_area("Text for Embeddings", height=150,
                                            placeholder="Enter text to generate embeddings for...",
                                            key="embed_text_input")
                with col2:
                    use_batch_emb = st.checkbox("Use Batch API", key="embed_batch")
                    st.info("Embeddings only use input tokens")
                
                if embed_text:
                    input_tokens = TokenCalculator.count_tokens(embed_text, model_id)
                    cost, breakdown = openai_models.calculate_embeddings_cost(model_id, input_tokens, use_batch_emb)
                    
                    col1, col2, col3, col4 = st.columns(4)
                    with col1:
                        st.metric("Input Tokens", f"{input_tokens:,}")
                    with col2:
                        st.metric("Dimensions", breakdown.get('dimensions', 'N/A'))
                    with col3:
                        st.metric("Cost/1K tokens", f"${breakdown['cost_per_1k_tokens']:.4f}")
                    with col4:
                        st.metric("Total Cost", f"${cost:.6f}")
        
        elif testing_category == "🎨 Image Generation":
            st.write("### Image Generation Testing")
            
            if openai_models.image_generation_models:
                model_id = st.selectbox("Select Image Model", list(openai_models.image_generation_models.keys()), key="img_gen_test_model")
                model_data = openai_models.image_generation_models[model_id]
                
                col1, col2, col3 = st.columns(3)
                with col1:
                    num_images = st.number_input("Number of Images", min_value=1, value=1, max_value=10, key="img_gen_count")
                with col2:
                    resolution = st.selectbox("Resolution", list(model_data['resolutions'].keys()), key="img_gen_resolution")
                with col3:
                    quality_options = list(model_data['resolutions'][resolution].keys())
                    quality = st.selectbox("Quality", quality_options, key="img_gen_quality")
                
                # Calculate cost
                cost, breakdown = openai_models.calculate_image_generation_cost(model_id, num_images, resolution, quality)
                
                # Results
                col1, col2, col3, col4 = st.columns(4)
                with col1:
                    st.metric("Images", num_images)
                with col2:
                    st.metric("Resolution", resolution)
                with col3:
                    st.metric("Cost/Image", f"${breakdown['cost_per_image']:.4f}")
                with col4:
                    st.metric("Total Cost", f"${cost:.4f}")
                
                # Sample prompts
                with st.expander("💡 Sample Image Prompts"):
                    sample_prompts = [
                        "A futuristic cityscape at sunset with flying cars",
                        "A detailed portrait of a wise old wizard",
                        "An abstract painting in the style of Van Gogh",
                        "A photorealistic image of a mountain landscape",
                        "A cute cartoon character for a children's book"
                    ]
                    for prompt in sample_prompts:
                        st.write(f"• {prompt}")
        
        elif testing_category == "🔍 Search & Tools":
            st.write("### Search and Tools Testing")
            
            tool_subcategory = st.radio("Tool Type", 
                                      ["🌐 Web Search", "🛠️ Built-in Tools"],
                                      key="tool_subcategory")
            
            if tool_subcategory == "🌐 Web Search":
                if openai_models.web_search_models:
                    model_id = st.selectbox("Select Search Model", list(openai_models.web_search_models.keys()), key="search_test_model")
                    model_data = openai_models.web_search_models[model_id]
                    
                    col1, col2 = st.columns(2)
                    with col1:
                        num_calls = st.number_input("Number of Search Calls", min_value=1, value=1000, step=100, key="search_calls")
                    with col2:
                        context_size = st.selectbox("Context Size", list(model_data['search_context_pricing'].keys()), key="search_context")
                    
                    # Calculate cost
                    cost, breakdown = openai_models.calculate_web_search_cost(model_id, num_calls, context_size)
                    
                    col1, col2, col3 = st.columns(3)
                    with col1:
                        st.metric("Search Calls", f"{num_calls:,}")
                    with col2:
                        st.metric("Cost/1K Calls", f"${breakdown['cost_per_1k_calls']:.2f}")
                    with col3:
                        st.metric("Total Cost", f"${cost:.4f}")
            
            elif tool_subcategory == "🛠️ Built-in Tools":
                if openai_models.built_in_tools:
                    tool_id = st.selectbox("Select Tool", list(openai_models.built_in_tools.keys()), key="builtin_test_tool")
                    tool_data = openai_models.built_in_tools[tool_id]
                    
                    if tool_id == 'code_interpreter':
                        containers = st.number_input("Number of Code Interpreter Sessions", min_value=1, value=1, key="code_containers")
                        cost, breakdown = openai_models.calculate_tool_cost(tool_id, num_containers=containers)
                        
                        col1, col2 = st.columns(2)
                        with col1:
                            st.metric("Sessions", containers)
                        with col2:
                            st.metric("Total Cost", f"${cost:.4f}")
                    
                    elif tool_id == 'file_search_storage':
                        gb_days = st.number_input("Storage (GB-Days)", min_value=0.0, value=1.0, step=0.1, key="storage_gb_days")
                        cost, breakdown = openai_models.calculate_tool_cost(tool_id, gb_days=gb_days)
                        
                        col1, col2, col3 = st.columns(3)
                        with col1:
                            st.metric("GB-Days", f"{gb_days:.1f}")
                        with col2:
                            st.metric("Free GB", breakdown['free_gb'])
                        with col3:
                            st.metric("Total Cost", f"${cost:.4f}")
                    
                    elif tool_id == 'file_search_tool_call':
                        calls = st.number_input("Number of File Search Calls", min_value=1, value=1000, key="file_search_calls")
                        cost, breakdown = openai_models.calculate_tool_cost(tool_id, num_calls=calls)
                        
                        col1, col2 = st.columns(2)
                        with col1:
                            st.metric("Calls", f"{calls:,}")
                        with col2:
                            st.metric("Total Cost", f"${cost:.4f}")
        
        elif testing_category == "🎓 Fine-tuning Models":
            st.write("### Fine-tuning Models Testing")
            
            finetuning_models = {k: v for k, v in openai_models.models.items() 
                               if v.get('pricing_type') == 'fine_tuning'}
            
            if finetuning_models:
                model_id = st.selectbox("Select Fine-tuning Model", list(finetuning_models.keys()), key="ft_test_model")
                model_info = finetuning_models[model_id]
                
                st.info("Fine-tuning involves training costs and inference costs")
                
                col1, col2 = st.columns(2)
                with col1:
                    st.write("**Training Phase:**")
                    if 'training_cost_per_hour' in model_info:
                        training_hours = st.number_input("Training Hours", min_value=0.0, value=1.0, step=0.1, key="ft_hours")
                        training_tokens = 0
                    else:
                        training_tokens = st.number_input("Training Tokens", min_value=0, value=1000000, step=100000, key="ft_tokens")
                        training_hours = 0
                
                with col2:
                    st.write("**Inference Phase:**")
                    inference_input = st.number_input("Inference Input Tokens", min_value=0, value=1000, step=100, key="ft_inf_input")
                    inference_output = st.number_input("Inference Output Tokens", min_value=0, value=100, step=10, key="ft_inf_output")
                
                use_cached_ft = st.checkbox("Use Cached Input", key="ft_cached")
                use_data_sharing = st.checkbox("Use Data Sharing (if available)", 
                                             disabled=not model_info.get('has_data_sharing'), key="ft_data_sharing")
                
                # Calculate cost
                cost, breakdown = openai_models.calculate_fine_tuning_cost(
                    model_id, training_tokens, training_hours, inference_input, inference_output,
                    use_cached_ft, use_data_sharing
                )
                
                # Results
                col1, col2, col3, col4 = st.columns(4)
                with col1:
                    if training_hours > 0:
                        st.metric("Training Hours", f"{training_hours:.1f}")
                    else:
                        st.metric("Training Tokens", f"{training_tokens:,}")
                with col2:
                    st.metric("Training Cost", f"${breakdown['training_cost']:.4f}")
                with col3:
                    st.metric("Inference Cost", f"${breakdown['inference_input_cost'] + breakdown['inference_output_cost']:.4f}")
                with col4:
                    st.metric("Total Cost", f"${cost:.4f}")
        
        elif testing_category == "🆓 Free Models":
            st.write("### Free Models Testing")
            
            free_models = {k: v for k, v in openai_models.models.items() 
                          if v.get('is_free', False)}
            
            if free_models:
                model_id = st.selectbox("Select Free Model", list(free_models.keys()), key="free_test_model")
                
                test_text = st.text_area("Test Input", height=150,
                                       placeholder="Enter text to test with free models (typically moderation)...",
                                       key="free_text_input")
                
                if test_text:
                    tokens = TokenCalculator.count_tokens(test_text, model_id)
                    cost, breakdown = openai_models.calculate_free_model_cost(model_id)
                    
                    col1, col2, col3 = st.columns(3)
                    with col1:
                        st.metric("Input Tokens", f"{tokens:,}")
                    with col2:
                        st.metric("Model Type", breakdown.get('pricing_type', 'Free'))
                    with col3:
                        st.metric("Cost", "FREE ✨")
                    
                    st.success("💡 This model is completely free to use!")
            else:
                st.info("No free models currently loaded.")
        
        # Quick testing tips
        with st.expander("💡 Quick Testing Tips"):
            st.markdown("""
            **Efficient Testing Strategies:**
            - **Text Models**: Start with short prompts to understand token consumption
            - **Audio Models**: Consider both text and audio token requirements
            - **Embeddings**: Perfect for semantic search and similarity tasks
            - **Image Generation**: Resolution and quality significantly impact costs
            - **Fine-tuning**: Training costs are one-time, inference costs are ongoing
            - **Batch API**: Use for non-real-time workloads to save 50%
            - **Cached Input**: Leverage for repeated content to save 50%
            """)
    
    with tab7:
        st.header("⚙️ Model Management")
        st.markdown("Manage your model data files - add, edit, delete, and maintain model configurations.")
        
        # Check if model manager is available
        if hasattr(openai_models, 'model_manager'):
            show_model_management_ui(openai_models.model_manager)
        else:
            st.error("Model Manager not available. Please check your configuration.")
    
    # Footer
    st.markdown("---")
    st.markdown("💡 **Tips:**")
    st.markdown("- Use cached input for repeated content to save 50%")
    st.markdown("- Consider Batch API for non-real-time workloads to save 50%")
    st.markdown("- Different models use different encodings for token counting")
    st.markdown("- Image generation and special tools have different pricing structures")

if __name__ == "__main__":
    main() 